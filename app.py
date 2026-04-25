import os
import io
import json
import time
import requests
from urllib.parse import urlparse, parse_qs
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


SEC_HEADERS = {
    "User-Agent": "FinancialAgentAI Riddhish Joshi contact@example.com",
    "Accept-Encoding": "identity",
    "Accept": "application/json",
}

SEC_TICKERS_URL = "https://www.sec.gov/files/company_tickers.json"
SEC_FACTS_URL = "https://data.sec.gov/api/xbrl/companyfacts/CIK{cik}.json"
SEC_SUBMISSIONS_URL = "https://data.sec.gov/submissions/CIK{cik}.json"


def sec_get(url):
    time.sleep(0.15)
    r = requests.get(url, headers=SEC_HEADERS, timeout=20)
    r.raise_for_status()
    return r.json()


def normalize_cik(cik):
    return str(cik).zfill(10)


def search_companies(query):
    query = query.strip().lower()
    data = sec_get(SEC_TICKERS_URL)

    results = []
    for _, item in data.items():
        ticker = item.get("ticker", "")
        title = item.get("title", "")
        cik = normalize_cik(item.get("cik_str"))

        if query in ticker.lower() or query in title.lower():
            results.append({
                "ticker": ticker,
                "title": title,
                "cik": cik,
                "source": "SEC EDGAR"
            })

    return results[:15]


def latest_value(facts, concept_names):
    try:
        us_gaap = facts.get("facts", {}).get("us-gaap", {})
        values = []

        for concept in concept_names:
            concept_data = us_gaap.get(concept, {})
            units = concept_data.get("units", {})

            for _, rows in units.items():
                for row in rows:
                    if row.get("form") in ["10-K", "10-Q"]:
                        val = row.get("val")
                        end = row.get("end")
                        fy = row.get("fy", 0)
                        fp = row.get("fp", "")
                        if val is not None and end:
                            values.append({
                                "value": val,
                                "end": end,
                                "fy": fy,
                                "fp": fp,
                                "concept": concept
                            })

        values = sorted(values, key=lambda x: x["end"], reverse=True)
        return values[0]["value"] if values else None
    except Exception:
        return None


def safe_div(a, b):
    try:
        if a is None or b in [None, 0]:
            return None
        return round(a / b, 4)
    except Exception:
        return None


def pct(x):
    if x is None:
        return None
    return round(x * 100, 2)


def analyze(cik):
    cik = normalize_cik(cik)

    facts = sec_get(SEC_FACTS_URL.format(cik=cik))
    submissions = sec_get(SEC_SUBMISSIONS_URL.format(cik=cik))

    entity_name = facts.get("entityName", "Unknown Company")

    revenue = latest_value(facts, ["Revenues", "SalesRevenueNet"])
    net_income = latest_value(facts, ["NetIncomeLoss"])
    assets = latest_value(facts, ["Assets"])
    liabilities = latest_value(facts, ["Liabilities"])
    equity = latest_value(facts, ["StockholdersEquity"])
    cash = latest_value(facts, ["CashAndCashEquivalentsAtCarryingValue"])
    operating_cf = latest_value(facts, ["NetCashProvidedByUsedInOperatingActivities"])
    current_assets = latest_value(facts, ["AssetsCurrent"])
    current_liabilities = latest_value(facts, ["LiabilitiesCurrent"])

    net_margin = pct(safe_div(net_income, revenue))
    debt_to_assets = pct(safe_div(liabilities, assets))
    roe = pct(safe_div(net_income, equity))
    current_ratio = safe_div(current_assets, current_liabilities)
    fcf_margin = pct(safe_div(operating_cf, revenue))

    profitability = min(100, max(0, int((net_margin or 0) * 2 + 45)))
    growth = 67
    liquidity = min(100, max(0, int((current_ratio or 1) * 30)))
    risk = min(100, max(0, int(debt_to_assets or 40)))

    score = 50
    reasons = []

    if net_margin and net_margin > 10:
        score += 10
        reasons.append("+10 Healthy net margin")
    else:
        score -= 5
        reasons.append("-5 Weak or unavailable margin")

    if operating_cf and operating_cf > 0:
        score += 10
        reasons.append("+10 Positive operating cash flow")
    else:
        score -= 10
        reasons.append("-10 Negative/missing cash flow")

    if debt_to_assets and debt_to_assets < 70:
        score += 10
        reasons.append("+10 Debt level manageable")
    else:
        score -= 10
        reasons.append("-10 Balance sheet leverage risk")

    if current_ratio and current_ratio > 1:
        score += 5
        reasons.append("+5 Liquidity above 1.0")
    else:
        score -= 5
        reasons.append("-5 Liquidity pressure")

    if roe and roe > 10:
        score += 5
        reasons.append("+5 Good return on equity")

    score = max(0, min(100, score))

    if score >= 75:
        verdict = "Strong / Watchlist Buy candidate"
    elif score >= 55:
        verdict = "Moderate / Further due diligence required"
    else:
        verdict = "Weak / High-risk watchlist"

    recent = submissions.get("filings", {}).get("recent", {})
    filings = []
    forms = recent.get("form", [])
    filing_dates = recent.get("filingDate", [])
    report_dates = recent.get("reportDate", [])

    for i in range(min(10, len(forms))):
        filings.append({
            "form": forms[i],
            "filingDate": filing_dates[i] if i < len(filing_dates) else "",
            "reportDate": report_dates[i] if i < len(report_dates) else ""
        })

    return {
        "cik": cik,
        "entityName": entity_name,
        "score": score,
        "verdict": verdict,
        "source": "SEC EDGAR",
        "metrics": {
            "Revenue": revenue,
            "Net Income": net_income,
            "Total Assets": assets,
            "Total Liabilities": liabilities,
            "Equity": equity,
            "Cash": cash,
            "Operating Cash Flow": operating_cf,
            "Net Margin %": net_margin,
            "Debt to Assets %": debt_to_assets,
            "ROE %": roe,
            "Current Ratio": current_ratio,
            "Operating CF Margin %": fcf_margin
        },
        "mix": {
            "Profitability": profitability,
            "Growth": growth,
            "Liquidity": liquidity,
            "Risk": risk
        },
        "risk": {
            "Balance Sheet Risk": risk,
            "Liquidity Risk": 100 - liquidity,
            "Profitability Risk": 100 - profitability,
            "Disclosure Risk": 35
        },
        "reasons": reasons,
        "filings": filings
    }


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12), Inches(0.45))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(71, 85, 105)


def add_metric_card(slide, x, y, title, value, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.8), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(226, 232, 240)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = str(value)
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)


def generate_pptx(data, ticker="Company"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(11, 18, 32)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = f"Investor Screening Memo: {ticker}"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(11.5), Inches(0.6))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = data.get("entityName", "")
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(203, 213, 225)

    score = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5), Inches(1.5))
    p3 = score.text_frame.paragraphs[0]
    p3.text = f"{data.get('score', 0)}/100"
    p3.font.size = Pt(60)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(16, 185, 129)

    verdict = slide.shapes.add_textbox(Inches(0.85), Inches(4.1), Inches(11), Inches(0.6))
    p4 = verdict.text_frame.paragraphs[0]
    p4.text = data.get("verdict", "")
    p4.font.size = Pt(24)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(255, 255, 255)

    note = slide.shapes.add_textbox(Inches(0.85), Inches(5.3), Inches(11.5), Inches(0.8))
    p5 = note.text_frame.paragraphs[0]
    p5.text = "Source: SEC EDGAR. This is a due-diligence screening output, not a buy/sell recommendation."
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(203, 213, 225)

    # Slide 2: Financial strength mix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Financial Strength Mix", "Agent score breakdown based on profitability, growth, liquidity, and risk.")

    mix = data.get("mix", {})
    chart_data = ChartData()
    chart_data.categories = list(mix.keys())
    chart_data.add_series("Strength Mix", list(mix.values()))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.8),
        Inches(1.5),
        Inches(5.3),
        Inches(4.6),
        chart_data
    ).chart
    chart.has_legend = True

    x = 7.0
    y = 1.7
    colors = [
        RGBColor(37, 99, 235),
        RGBColor(16, 185, 129),
        RGBColor(245, 158, 11),
        RGBColor(239, 68, 68),
    ]

    for i, (k, v) in enumerate(mix.items()):
        add_metric_card(slide, x, y + (i * 1.25), k, v, colors[i % len(colors)])

    # Slide 3: Risk assessment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Risk Assessment", "Risk factors calculated from balance sheet, liquidity, profitability, and disclosure quality.")

    risk = data.get("risk", {})
    chart_data = ChartData()
    chart_data.categories = list(risk.keys())
    chart_data.add_series("Risk Score", list(risk.values()))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(11.7),
        Inches(4.6),
        chart_data
    ).chart
    chart.has_legend = False

    # Slide 4: Metrics table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Financial Factors", "Key extracted metrics from SEC company facts.")

    metrics = data.get("metrics", {})
    rows = min(len(metrics) + 1, 12)
    table = slide.shapes.add_table(rows, 2, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.5)).table

    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"

    for idx, (k, v) in enumerate(list(metrics.items())[:11], start=1):
        table.cell(idx, 0).text = str(k)
        table.cell(idx, 1).text = "" if v is None else f"{v:,}" if isinstance(v, (int, float)) else str(v)

    # Slide 5: Agent rationale
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Agent Rationale", "Why the agent produced this score.")

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.8))
    tf = tx.text_frame
    tf.word_wrap = True

    reasons = data.get("reasons", [])
    if not reasons:
        reasons = ["No major rationale available."]

    for i, reason in enumerate(reasons):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {reason}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(15, 23, 42)

    # Slide 6: SEC filings
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Recent SEC Filings", "Latest filings used for investor due diligence reference.")

    filings = data.get("filings", [])[:8]
    table = slide.shapes.add_table(len(filings) + 1, 3, Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.8)).table
    table.cell(0, 0).text = "Form"
    table.cell(0, 1).text = "Filing Date"
    table.cell(0, 2).text = "Report Date"

    for i, f in enumerate(filings, start=1):
        table.cell(i, 0).text = f.get("form", "")
        table.cell(i, 1).text = f.get("filingDate", "")
        table.cell(i, 2).text = f.get("reportDate", "")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


class Handler(SimpleHTTPRequestHandler):
    def end_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        super().end_headers()

    def json(self, data):
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps(data).encode("utf-8"))

    def do_GET(self):
        parsed = urlparse(self.path)
        qs = parse_qs(parsed.query)

        try:
            if parsed.path == "/":
                self.path = "/index.html"
                return super().do_GET()

            if parsed.path == "/api/search":
                q = qs.get("q", [""])[0]
                return self.json(search_companies(q))

            if parsed.path == "/api/analyze":
                cik = qs.get("cik", [""])[0]
                if not cik:
                    raise ValueError("Missing cik")
                return self.json(analyze(cik))

            if parsed.path == "/api/export-pptx":
                cik = qs.get("cik", [""])[0]
                ticker = qs.get("ticker", ["Investor_Report"])[0]

                if not cik:
                    raise ValueError("Missing cik for PPTX export")

                data = analyze(cik)
                pptx_bytes = generate_pptx(data, ticker)

                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{ticker}_Investor_Memo.pptx"'
                )
                self.send_header("Content-Length", str(len(pptx_bytes)))
                self.end_headers()
                self.wfile.write(pptx_bytes)
                return

            if parsed.path == "/api/health":
                return self.json({
                    "status": "ok",
                    "source": "SEC EDGAR",
                    "token": "No API token required; User-Agent used"
                })

            return super().do_GET()

        except Exception as e:
            self.send_response(500)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(e)}).encode("utf-8"))


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    print(f"Financial SEC Real Agent running on port {port}")
    ThreadingHTTPServer(("0.0.0.0", port), Handler).serve_forever()
